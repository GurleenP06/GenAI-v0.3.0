import logging
from pathlib import Path

import pptx as pptx_lib

logger = logging.getLogger(__name__)


def extract_text_from_pptx(pptx_path: Path) -> str:
    try:
        presentation = pptx_lib.Presentation(str(pptx_path))
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting PPTX {pptx_path}: {e}")
        return ""
